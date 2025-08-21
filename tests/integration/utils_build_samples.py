# Copyright (c) 2025 takotime808

from __future__ import annotations

import fitz
import openpyxl
import pandas as pd
from PIL import Image, ImageDraw
from pathlib import Path
from docx import Document
from pptx import Presentation
from pptx.util import Inches

def build_txt(path: Path) -> Path:
    path.write_text("hello\nworld\n")
    return path

def build_html(path: Path) -> Path:
    path.write_text("<html><head><title>T</title></head><body><h1>Header</h1><p>Hello <b>world</b>!</p></body></html>")
    return path

def build_docx(path: Path) -> Path:
    doc = Document()
    doc.add_paragraph("Paragraph one")
    doc.add_paragraph("Paragraph two")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "A"
    table.rows[0].cells[1].text = "B"
    table.rows[1].cells[0].text = "C"
    table.rows[1].cells[1].text = "D"
    doc.save(path)
    return path

def build_pptx(path: Path) -> Path:
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(slide_layout)
    left = top = Inches(1)
    width = Inches(8)
    height = Inches(1.5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tx.text_frame.text = "Hello slide"
    prs.save(path)
    return path

def build_xlsx(path: Path) -> Path:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["a","b"])
    ws.append([1,2])
    wb.save(path)
    return path

def build_csv(path: Path) -> Path:
    df = pd.DataFrame({"x":[1,2],"y":[3,4]})
    df.to_csv(path, index=False)
    return path

def build_image_png(path: Path) -> Path:
    img = Image.new("RGB", (220, 90), (255,255,255))
    d = ImageDraw.Draw(img)
    d.text((10, 35), "HELLO", fill=(0,0,0))
    img.save(path)
    return path

def build_pdf(path: Path) -> Path:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72,72), "Hello PDF Page 1")
    doc.save(str(path))
    doc.close()
    return path
