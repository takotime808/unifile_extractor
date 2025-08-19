# Copyright (c) 2025 takotime808

import pytest
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from unifile.extractors.pptx_extractor import PptxExtractor

def _build_pptx(path: Path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # title-only / blank
    slide = prs.slides.add_slide(slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Hello slide"
    prs.save(path)

def test_pptx_extractor_reads_slide_text(tmp_path):
    p = tmp_path / "sample.pptx"
    _build_pptx(p)
    ext = PptxExtractor()
    rows = ext.extract(p)
    assert rows, "should have at least one slide row"
    r0 = rows[0]
    assert r0.unit_type == "slide"
    assert "Hello slide" in r0.content
