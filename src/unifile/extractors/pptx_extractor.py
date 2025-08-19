# Copyright (c) 2025 takotime808

from __future__ import annotations

from typing import List
from pathlib import Path
from pptx import Presentation

from unifile.extractors.base import (
    BaseExtractor,
    make_row,
    Row,
)


class PptxExtractor(BaseExtractor):
    """
    PPTX --> plain-text extractor.

    For each slide, this extractor gathers text from shapes that expose a
    `.text` attribute and emits one standardized row per slide.

    Inherits :meth:`BaseExtractor.extract` for path validation and
    exception-to-error-row wrapping. The slide parsing logic is implemented
    in :meth:`_extract`.

    Output Rows
    -----------
    * Per slide (index `i`):
        - file_type: "pptx"
        - unit_type: "slide"
        - unit_id:   str(i)
        - content:   Combined newline-joined text from slide shapes
        - metadata:  {"slide_index": i}

    Notes
    -----
    - Shapes without textual content are skipped.
    - Errors on individual shapes are ignored, so one bad shape does not
      prevent the rest of the slide from being processed.
    """

    supported_extensions = ["pptx"]

    def _extract(self, path: Path) -> List[Row]:
        """
        Parse a PPTX file into standardized per-slide rows.

        Parameters
        ----------
        path
            Path to a `.pptx` file. Existence checks are handled by
            :class:`BaseExtractor.extract`.

        Returns
        -------
        list[Row]
            One row per slide with collected text content.
        """
        prs = Presentation(str(path))
        rows: List[Row] = []

        for i, slide in enumerate(prs.slides):
            texts: list[str] = []
            for shape in slide.shapes:
                # Some shapes (e.g., pictures, charts, tables) may not have .text;
                # guard with hasattr and swallow unexpected shape-level errors.
                try:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                except Exception:
                    # Ignore problematic shapes; proceed with the rest
                    continue

            content = "\n".join(texts).strip()
            rows.append(
                make_row(
                    path=path,
                    file_type="pptx",
                    unit_type="slide",
                    unit_id=str(i),
                    content=content,
                    metadata={"slide_index": i},
                    status="ok",
                )
            )

        return rows
